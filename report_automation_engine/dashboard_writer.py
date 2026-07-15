"""Create portrait organization dashboard PPTX files from dashboard_package.json."""

from __future__ import annotations

import argparse
import json
import tempfile
from pathlib import Path
from typing import Any, Dict, List

from copy import deepcopy

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


BLUE = RGBColor(0x2F, 0x66, 0xA7)
LIGHT_BLUE = RGBColor(0xE9, 0xF1, 0xF8)
ACCENT = RGBColor(0xD9, 0x42, 0x2E)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF3, 0xF5, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

THEMES = {
    "modern_blue": {
        "background": RGBColor(0xD8, 0xF0, 0xF5),
        "card": WHITE,
        "card_alt": RGBColor(0xEC, 0xF3, 0xFB),
        "border": RGBColor(0xD8, 0xE2, 0xEC),
        "primary": RGBColor(0x21, 0x4E, 0x7A),
        "pill": RGBColor(0x0E, 0x35, 0x4F),
        "secondary": RGBColor(0x5F, 0x72, 0x85),
        "accent": RGBColor(0xF0, 0x78, 0x4A),
        "muted": RGBColor(0xF0, 0xF3, 0xF6),
    },
    "modern_mint": {
        "background": RGBColor(0xF5, 0xFA, 0xF8),
        "card": WHITE,
        "card_alt": RGBColor(0xE8, 0xF6, 0xF0),
        "border": RGBColor(0xC9, 0xE4, 0xDA),
        "primary": RGBColor(0x1B, 0x68, 0x5B),
        "pill": RGBColor(0x10, 0x48, 0x42),
        "secondary": RGBColor(0x5E, 0x73, 0x6D),
        "accent": RGBColor(0xE8, 0x76, 0x55),
        "muted": RGBColor(0xEC, 0xF2, 0xEF),
    },
    "graphite": {
        "background": RGBColor(0xF7, 0xF7, 0xF5),
        "card": WHITE,
        "card_alt": RGBColor(0xF0, 0xF0, 0xEC),
        "border": RGBColor(0xD9, 0xD9, 0xD1),
        "primary": RGBColor(0x32, 0x37, 0x3D),
        "pill": RGBColor(0x29, 0x2E, 0x34),
        "secondary": RGBColor(0x6C, 0x70, 0x76),
        "accent": RGBColor(0xC9, 0x69, 0x46),
        "muted": RGBColor(0xED, 0xED, 0xE8),
    },
}


def write_dashboard(
    package_path: str | Path,
    preflight_path: str | Path,
    output_path: str | Path,
    template_path: str | Path | None = None,
) -> Path:
    package = load_json(package_path)
    preflight = load_json(preflight_path)
    if preflight.get("status") == "blocked":
        raise ValueError("dashboard preflight is blocked; fix errors before creating PPTX output")

    theme = dashboard_theme(package)
    if template_path and Path(template_path).exists():
        prs = build_from_template(package, template_path, theme)
    else:
        prs = Presentation()
        set_page_size(prs, str(package.get("meta", {}).get("page_size", "A4")))
        remove_default_slides(prs)
        for entity in package.get("entities", []):
            add_dashboard_slide(prs, package, entity, theme)

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return output


def build_from_template(package: Dict[str, Any], template_path: str | Path, theme: Dict[str, Any]) -> Presentation:
    prs = Presentation(str(template_path))
    if not prs.slides:
        raise ValueError("dashboard template has no slides")

    template_slide = prs.slides[0]
    frame = extract_template_frame(template_slide, theme)
    entities = package.get("entities", [])
    while len(prs.slides) < len(entities):
        clone_slide(prs, template_slide)
    while len(prs.slides) > len(entities):
        remove_slide(prs, len(prs.slides) - 1)

    for idx, entity in enumerate(entities):
        apply_template_slide(prs.slides[idx], package, entity, frame, theme, int(prs.slide_width), int(prs.slide_height))
    return prs


def set_page_size(prs: Presentation, page_size: str) -> None:
    if page_size.upper() == "B5":
        prs.slide_width = Inches(6.93)
        prs.slide_height = Inches(9.84)
    else:
        prs.slide_width = Inches(8.27)
        prs.slide_height = Inches(11.69)


def remove_default_slides(prs: Presentation) -> None:
    for slide_id in list(prs.slides._sldIdLst):
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)


def remove_slide(prs: Presentation, index: int) -> None:
    slide_id = list(prs.slides._sldIdLst)[index]
    prs.part.drop_rel(slide_id.rId)
    prs.slides._sldIdLst.remove(slide_id)


def clone_slide(prs: Presentation, source_slide):
    new_slide = prs.slides.add_slide(blank_layout(prs))
    rel_map = {}
    for rel in source_slide.part.rels.values():
        if "notesSlide" not in rel.reltype:
            rel_map[rel.rId] = new_slide.part.rels._add_relationship(rel.reltype, rel._target, rel.is_external)
    for shape in source_slide.shapes:
        element = deepcopy(shape.element)
        replace_relationship_ids(element, rel_map)
        new_slide.shapes._spTree.insert_element_before(element, "p:extLst")
    return new_slide


def replace_relationship_ids(element, rel_map: Dict[str, str]) -> None:
    for node in [element, *list(element.iter())]:
        for attr, value in list(node.attrib.items()):
            if value in rel_map:
                node.set(attr, rel_map[value])


def extract_template_frame(slide, theme: Dict[str, Any]) -> Dict[str, Any]:
    shapes = list(slide.shapes)
    by_name = {shape.name: shape for shape in shapes}
    frame: Dict[str, Any] = {"font": theme["font"], "texts": {}, "slots": {}}
    for shape in shapes:
        if getattr(shape, "has_text_frame", False):
            font = first_font_name(shape)
            if font:
                frame["font"] = font
                break
    for key in ["RA_DASH_TITLE", "RA_DASH_PROFILE", "RA_DASH_SOURCE"]:
        if key in by_name:
            frame["texts"][key] = style_box(by_name[key])
    frame["texts"]["RA_DASH_NARRATIVE"] = child_text_box(shapes, by_name.get("RA_DASH_NARRATIVE")) or (
        style_box(by_name["RA_DASH_NARRATIVE"]) if "RA_DASH_NARRATIVE" in by_name else None
    )
    if "RA_DASH_NARRATIVE_TEXT" in by_name:
        frame["texts"]["RA_DASH_NARRATIVE"] = style_box(by_name["RA_DASH_NARRATIVE_TEXT"])
    for idx in range(1, 7):
        card = by_name.get(f"RA_DASH_KPI_{idx}")
        children = child_text_boxes(shapes, card)
        label = by_name.get(f"RA_DASH_KPI_{idx}_LABEL")
        value = by_name.get(f"RA_DASH_KPI_{idx}_VALUE")
        frame["slots"][f"kpi_{idx}"] = {
            "card": style_box(card) if card else None,
            "label": style_box(label) if label else (style_box(children[0]) if len(children) > 0 else None),
            "value": style_box(value) if value else (style_box(children[1]) if len(children) > 1 else None),
        }
    for idx in range(1, 5):
        card = by_name.get(f"RA_DASH_CHART_{idx}")
        children = child_text_boxes(shapes, card)
        chart = child_chart_box(shapes, card)
        title = by_name.get(f"RA_DASH_CHART_{idx}_TITLE")
        plot = by_name.get(f"RA_DASH_CHART_{idx}_PLOT")
        frame["slots"][f"chart_{idx}"] = {
            "card": style_box(card) if card else None,
            "title": style_box(title) if title else (style_box(children[0]) if children else None),
            "chart": style_box(plot) if plot else (style_box(chart) if chart else chart_box_from_card(card)),
        }
    return frame


def apply_template_slide(slide, package: Dict[str, Any], entity: Dict[str, Any], frame: Dict[str, Any], theme: Dict[str, Any], slide_width: int, slide_height: int) -> None:
    remove_dynamic_shapes(slide)
    set_text(slide, "RA_DASH_TITLE", entity.get("entity_name", ""), frame, theme, bold=True)
    profile = "   |   ".join(f"{item.get('label')}: {item.get('value')}" for item in entity.get("profile", []) if item.get("value"))
    set_text(slide, "RA_DASH_PROFILE", profile, frame, theme)
    set_text(slide, "RA_DASH_SOURCE", source_text(package, entity), frame, theme)
    add_template_text(slide, frame["texts"].get("RA_DASH_NARRATIVE"), entity.get("narrative", ""), frame, theme, "RA_DASH_NARRATIVE_TEXT")

    for idx, kpi in enumerate(entity.get("kpis", [])[:6], start=1):
        slot = frame["slots"].get(f"kpi_{idx}", {})
        add_template_text(slide, slot.get("label"), str(kpi.get("label", "")), frame, theme, f"RA_DASH_KPI_{idx}_LABEL")
        add_template_text(slide, slot.get("value"), str(kpi.get("display_value", "")), frame, theme, f"RA_DASH_KPI_{idx}_VALUE", bold=True)

    for idx, chart in enumerate(entity.get("charts", [])[:4], start=1):
        slot = frame["slots"].get(f"chart_{idx}", {})
        add_template_text(slide, slot.get("title"), str(chart.get("title", "")), frame, theme, f"RA_DASH_CHART_{idx}_TITLE", bold=True)
        box = slot.get("chart")
        if box:
            add_chart(
                slide,
                chart.get("points", []),
                choose_chart_type(chart),
                box["left"],
                box["top"],
                box["width"],
                box["height"],
                theme,
            )
    clamp_slide_shapes(slide, slide_width, slide_height)


def clamp_slide_shapes(slide, width: int, height: int) -> None:
    for shape in slide.shapes:
        try:
            if int(shape.left) < 0:
                shape.left = 0
            if int(shape.top) < 0:
                shape.top = 0
            if int(shape.left + shape.width) > width:
                shape.width = max(0, width - int(shape.left))
            if int(shape.top + shape.height) > height:
                shape.height = max(0, height - int(shape.top))
        except Exception:
            continue


def remove_dynamic_shapes(slide) -> None:
    shapes = list(slide.shapes)
    card_names = {f"RA_DASH_KPI_{idx}" for idx in range(1, 7)} | {f"RA_DASH_CHART_{idx}" for idx in range(1, 5)}
    cards = [shape for shape in shapes if getattr(shape, "name", "") in card_names]
    narrative = next((shape for shape in shapes if getattr(shape, "name", "") == "RA_DASH_NARRATIVE"), None)
    if narrative is not None:
        cards.append(narrative)

    remove = []
    for shape in shapes:
        name = getattr(shape, "name", "")
        is_generated_text = (
            name.endswith("_LABEL")
            or name.endswith("_VALUE")
            or name.endswith("_TITLE")
            or name == "RA_DASH_NARRATIVE_TEXT"
            or name.endswith("_PLOT")
        )
        is_template_child_text = shape not in cards and getattr(shape, "has_text_frame", False) and any(inside(shape, card) for card in cards)
        if getattr(shape, "has_chart", False) or is_generated_text or is_template_child_text:
            remove.append(shape)
    for shape in remove:
        shape.element.getparent().remove(shape.element)


def set_text(slide, name: str, text: str, frame: Dict[str, Any], theme: Dict[str, Any], bold: bool = False) -> None:
    for shape in slide.shapes:
        if shape.name == name and getattr(shape, "has_text_frame", False):
            write_shape_text(shape, text, frame["font"], bold)
            return
    add_template_text(slide, frame["texts"].get(name), text, frame, theme, name, bold=bold)


def add_template_text(slide, box: Dict[str, Any] | None, text: str, frame: Dict[str, Any], theme: Dict[str, Any], name: str, bold: bool = False):
    if not box:
        return None
    shape = add_text(slide, text, box["left"], box["top"], box["width"], box["height"], box.get("size", 9), box.get("color") or theme["primary"], bold=bold or box.get("bold", False), name=name, theme={**theme, "font": box.get("font") or frame["font"]})
    return shape


def write_shape_text(shape, text: str, font: str, bold: bool = False) -> None:
    size = first_font_size(shape)
    color = first_font_color(shape)
    shape.text_frame.clear()
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = str(text)
    run.font.name = font
    if size:
        run.font.size = size
    run.font.bold = bold or first_font_bold(shape)
    if color:
        run.font.color.rgb = color


def style_box(shape) -> Dict[str, Any] | None:
    if shape is None:
        return None
    box = {
        "left": int(shape.left),
        "top": int(shape.top),
        "width": int(shape.width),
        "height": int(shape.height),
        "font": "",
        "size": 9,
        "color": None,
        "bold": False,
    }
    if getattr(shape, "has_text_frame", False):
        box["font"] = first_font_name(shape)
        box["size"] = first_font_size_pt(shape)
        box["color"] = first_font_color(shape)
        box["bold"] = first_font_bold(shape)
    return box


def child_text_boxes(shapes, parent) -> List[Any]:
    if parent is None:
        return []
    children = []
    for shape in shapes:
        if shape is parent or not getattr(shape, "has_text_frame", False):
            continue
        if inside(shape, parent):
            children.append(shape)
    children.sort(key=lambda item: (int(item.top), int(item.left)))
    return children


def child_text_box(shapes, parent):
    boxes = child_text_boxes(shapes, parent)
    return style_box(boxes[0]) if boxes else None


def child_chart_box(shapes, parent):
    if parent is None:
        return None
    for shape in shapes:
        if getattr(shape, "has_chart", False) and inside(shape, parent):
            return shape
    return None


def chart_box_from_card(card) -> Dict[str, Any] | None:
    if card is None:
        return None
    return {
        "left": int(card.left + Inches(0.18)),
        "top": int(card.top + Inches(0.45)),
        "width": int(card.width - Inches(0.36)),
        "height": int(card.height - Inches(0.62)),
    }


def inside(shape, parent) -> bool:
    return (
        int(shape.left) >= int(parent.left)
        and int(shape.top) >= int(parent.top)
        and int(shape.left + shape.width) <= int(parent.left + parent.width)
        and int(shape.top + shape.height) <= int(parent.top + parent.height)
    )


def first_font_name(shape) -> str:
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.name:
                return run.font.name
    return ""


def first_font_size(shape):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                return run.font.size
    return None


def first_font_size_pt(shape) -> int:
    size = first_font_size(shape)
    return int(round(size.pt)) if size else 9


def first_font_color(shape):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            try:
                if run.font.color and run.font.color.rgb:
                    return run.font.color.rgb
            except Exception:
                pass
    return None


def first_font_bold(shape) -> bool:
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.bold:
                return True
    return False


def source_text(package: Dict[str, Any], entity: Dict[str, Any]) -> str:
    source = package.get("meta", {}).get("source_file_name", "")
    created = package.get("meta", {}).get("created_at", "")
    return f"source: {source} / entity: {entity.get('entity_name', '')} / generated: {created}"


def blank_layout(prs: Presentation):
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]


def add_dashboard_slide(prs: Presentation, package: Dict[str, Any], entity: Dict[str, Any], theme: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    width = prs.slide_width
    height = prs.slide_height
    margin = Inches(0.45)
    footer_h = Inches(1.2)
    footer_top = height - footer_h
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = theme["background"]

    add_illustration_footer(slide, width, footer_top, footer_h, theme)
    add_header(slide, entity, margin, Inches(0.45), width - margin * 2, Inches(0.55), theme)
    add_profile(slide, entity, width - Inches(2.95), Inches(0.56), Inches(2.5), Inches(0.32), theme)

    add_section_label(slide, "기업/기관 핵심 지표", margin, Inches(1.18), Inches(2.2), theme)
    add_panel(slide, margin, Inches(1.55), width - margin * 2, Inches(2.25), "RA_DASH_PANEL_KPI", theme)
    add_kpi_grid(slide, entity.get("kpis", []), margin + Inches(0.22), Inches(1.8), width - margin * 2 - Inches(0.44), Inches(1.55), theme)
    add_narrative(slide, entity, margin + Inches(0.22), Inches(3.32), width - margin * 2 - Inches(0.44), Inches(0.34), theme)

    add_section_label(slide, "성과 및 인식 지표", margin, Inches(4.08), Inches(2.0), theme)
    chart_panel_top = Inches(4.45)
    chart_panel_h = footer_top - chart_panel_top - Inches(0.18)
    add_panel(slide, margin, chart_panel_top, width - margin * 2, chart_panel_h, "RA_DASH_PANEL_CHARTS", theme)
    add_charts(slide, entity.get("charts", []), margin + Inches(0.22), chart_panel_top + Inches(0.26), width - margin * 2 - Inches(0.44), chart_panel_h - Inches(0.5), theme)
    add_source(slide, package, entity, margin, height - Inches(0.22), width - margin * 2, Inches(0.18), theme)


def add_section_label(slide, text: str, left, top, width, theme: Dict[str, Any]) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(left), emu(top), emu(width), emu(Inches(0.28)))
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["pill"]
    shape.line.color.rgb = theme["pill"]
    add_text(slide, text, left + Inches(0.14), top + Inches(0.035), width - Inches(0.28), Inches(0.18), 8, WHITE, bold=True, theme=theme)


def add_panel(slide, left, top, width, height, name: str, theme: Dict[str, Any]) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(left), emu(top), emu(width), emu(height))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["card"]
    shape.line.color.rgb = theme["card"]


def add_header(slide, entity: Dict[str, Any], left, top, width, height, theme: Dict[str, Any]) -> None:
    title = str(entity.get("entity_name", "기관명"))
    pill_w = min(width, Inches(max(3.1, min(6.2, 1.2 + len(title) * 0.16))))
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(left), emu(top), emu(pill_w), emu(height))
    shape.name = "RA_DASH_TITLE_PILL"
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["pill"]
    shape.line.color.rgb = theme["pill"]
    add_text(slide, title, left + Inches(0.2), top + Inches(0.1), pill_w - Inches(0.4), height - Inches(0.16), 15, WHITE, bold=True, name="RA_DASH_TITLE", theme=theme)


def add_profile(slide, entity: Dict[str, Any], left, top, width, height, theme: Dict[str, Any]) -> None:
    profile = entity.get("profile", [])
    text = "   |   ".join(f"{item.get('label')}: {item.get('value')}" for item in profile if item.get("value"))
    add_text(slide, text or "프로필 정보 없음", left, top, width, height, 8, theme["secondary"], name="RA_DASH_PROFILE", theme=theme)


def add_illustration_footer(slide, slide_width, top, height, theme: Dict[str, Any]) -> None:
    try:
        path = create_footer_illustration(theme)
        slide.shapes.add_picture(str(path), 0, emu(top), width=emu(slide_width), height=emu(height))
    except Exception:
        ground = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, emu(top + height - Inches(0.18)), emu(slide_width), emu(Inches(0.18)))
        ground.fill.solid()
        ground.fill.fore_color.rgb = theme["primary"]
        ground.line.color.rgb = theme["primary"]


def add_kpi_grid(slide, kpis: List[Dict[str, Any]], left, top, width, height, theme: Dict[str, Any]) -> None:
    cols = 2
    rows = 3
    gap = Inches(0.12)
    card_w = (width - gap) / cols
    card_h = (height - gap * 2) / rows
    for idx in range(6):
        col = idx % cols
        row = idx // cols
        x = left + col * (card_w + gap)
        y = top + row * (card_h + gap)
        kpi = kpis[idx] if idx < len(kpis) else {}
        add_card(slide, x, y, card_w, card_h, f"RA_DASH_KPI_{idx + 1}", theme)
        add_text(slide, str(kpi.get("label", "")), x + Inches(0.14), y + Inches(0.09), card_w - Inches(0.28), Inches(0.20), 7, theme["secondary"], name=f"RA_DASH_KPI_{idx + 1}_LABEL", theme=theme)
        add_text(slide, str(kpi.get("display_value", "")), x + Inches(0.14), y + Inches(0.31), card_w - Inches(0.28), Inches(0.32), 14, theme["primary"], bold=True, name=f"RA_DASH_KPI_{idx + 1}_VALUE", theme=theme)


def add_narrative(slide, entity: Dict[str, Any], left, top, width, height, theme: Dict[str, Any]) -> None:
    add_card(slide, left, top, width, height, "RA_DASH_NARRATIVE", theme, alt=True)
    add_text(slide, str(entity.get("narrative", "")), left + Inches(0.16), top + Inches(0.13), width - Inches(0.32), height - Inches(0.2), 9, theme["primary"], name="RA_DASH_NARRATIVE_TEXT", theme=theme)


def add_charts(slide, charts: List[Dict[str, Any]], left, top, width, height, theme: Dict[str, Any]) -> None:
    cols = 2
    rows = 2
    gap = Inches(0.16)
    card_w = (width - gap) / cols
    card_h = (height - gap) / rows
    for idx in range(4):
        col = idx % cols
        row = idx // cols
        x = left + col * (card_w + gap)
        y = top + row * (card_h + gap)
        chart = charts[idx] if idx < len(charts) else {}
        add_card(slide, x, y, card_w, card_h, f"RA_DASH_CHART_{idx + 1}", theme)
        add_text(slide, str(chart.get("title", "")), x + Inches(0.14), y + Inches(0.11), card_w - Inches(0.28), Inches(0.22), 8, theme["primary"], bold=True, name=f"RA_DASH_CHART_{idx + 1}_TITLE", theme=theme)
        points = chart.get("points", [])
        chart_type = choose_chart_type(chart)
        if chart_type == "progress":
            add_progress(slide, points, x + Inches(0.25), y + Inches(0.55), card_w - Inches(0.5), card_h - Inches(0.85), theme)
        else:
            add_chart(slide, points, chart_type, x + Inches(0.18), y + Inches(0.45), card_w - Inches(0.36), card_h - Inches(0.62), theme)


def add_source(slide, package: Dict[str, Any], entity: Dict[str, Any], left, top, width, height, theme: Dict[str, Any]) -> None:
    source = package.get("meta", {}).get("source_file_name", "")
    created = package.get("meta", {}).get("created_at", "")
    add_text(slide, f"source: {source} / entity: {entity.get('entity_name', '')} / generated: {created}", left, top, width, height, 6, theme["secondary"], name="RA_DASH_SOURCE", theme=theme)


def add_chart(slide, points: List[Dict[str, Any]], chart_kind: str, left, top, width, height, theme: Dict[str, Any]) -> None:
    points = [point for point in points if number(point.get("value")) is not None]
    if not points:
        add_text(slide, "차트 데이터 없음", left, top, width, height, 10, theme["secondary"], theme=theme)
        return
    chart_data = ChartData()
    chart_data.categories = [str(point.get("category", "")) for point in points]
    chart_data.add_series("값", [number(point.get("value")) or 0 for point in points])
    if chart_kind == "line":
        ppt_chart_type = XL_CHART_TYPE.LINE_MARKERS
    elif chart_kind == "pie":
        ppt_chart_type = XL_CHART_TYPE.PIE
    else:
        ppt_chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    chart = slide.shapes.add_chart(ppt_chart_type, emu(left), emu(top), emu(width), emu(height), chart_data).chart
    chart.has_title = False
    chart.has_legend = ppt_chart_type == XL_CHART_TYPE.PIE
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
    highlight_max(chart, points, theme)


def add_progress(slide, points: List[Dict[str, Any]], left, top, width, height, theme: Dict[str, Any]) -> None:
    if not points:
        add_text(slide, "점수 데이터 없음", left, top, width, height, 10, theme["secondary"], theme=theme)
        return
    row_h = height / max(1, len(points))
    for idx, point in enumerate(points[:4]):
        value = max(0, min(100, number(point.get("value")) or 0))
        y = top + idx * row_h
        add_text(slide, str(point.get("category", "")), left, y, Inches(1.0), Inches(0.25), 7, theme["secondary"], theme=theme)
        track = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            emu(left + Inches(1.05)),
            emu(y),
            emu(width - Inches(1.45)),
            emu(Inches(0.18)),
        )
        track.fill.solid()
        track.fill.fore_color.rgb = theme["muted"]
        track.line.color.rgb = theme["muted"]
        fill_w = (width - Inches(1.45)) * (value / 100)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            emu(left + Inches(1.05)),
            emu(y),
            emu(fill_w),
            emu(Inches(0.18)),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = theme["accent"] if value == max(number(p.get("value")) or 0 for p in points) else theme["primary"]
        bar.line.color.rgb = bar.fill.fore_color.rgb
        add_text(slide, f"{value:.1f}", left + width - Inches(0.35), y - Inches(0.03), Inches(0.35), Inches(0.25), 6, theme["secondary"], theme=theme)


def choose_chart_type(chart: Dict[str, Any]) -> str:
    requested = str(chart.get("chart_type", "auto")).lower()
    points = chart.get("points", [])
    if requested in {"pie", "column", "line", "progress"}:
        return requested
    if len(points) == 1:
        return "progress"
    if len(points) == 2:
        return "pie"
    return "column"


def dashboard_theme(package: Dict[str, Any]) -> Dict[str, Any]:
    mapping = package.get("mapping", {}) if isinstance(package, dict) else {}
    preset = str(mapping.get("style_preset") or "modern_blue").strip()
    theme = dict(THEMES.get(preset, THEMES["modern_blue"]))
    theme["font"] = str(mapping.get("font_family") or "Malgun Gothic").strip() or "Malgun Gothic"
    return theme


def create_footer_illustration(theme: Dict[str, Any]) -> Path:
    from PIL import Image, ImageDraw

    key = f"{theme['primary']}_{theme['accent']}_{theme['background']}".replace(" ", "_").replace(",", "_")
    path = Path(tempfile.gettempdir()) / f"ra_dashboard_footer_{key}.png"
    if path.exists():
        return path

    w, h = 1600, 280
    img = Image.new("RGBA", (w, h), (255, 255, 255, 0))
    draw = ImageDraw.Draw(img)
    primary = rgb_tuple(theme["primary"])
    accent = rgb_tuple(theme["accent"])
    border = rgb_tuple(theme["border"])
    muted = rgb_tuple(theme["muted"])
    secondary = rgb_tuple(theme["secondary"])

    ground_y = 230
    draw.rectangle([0, ground_y, w, h], fill=(*primary, 255))

    buildings = [
        (760, 126, 86, 104, muted),
        (850, 96, 96, 134, border),
        (948, 116, 72, 114, secondary),
        (1026, 78, 130, 152, muted),
        (1168, 104, 80, 126, border),
        (1254, 72, 150, 158, secondary),
        (1416, 95, 95, 135, border),
    ]
    for x, y, bw, bh, color in buildings:
        draw.rectangle([x, y, x + bw, ground_y], fill=(*color, 230))
        for wx in range(x + 14, x + bw - 8, 26):
            for wy in range(y + 18, min(ground_y - 15, y + bh), 28):
                draw.rectangle([wx, wy, wx + 10, wy + 12], fill=(255, 255, 255, 145))

    for x in (70, 580, 690, 1500):
        draw.rectangle([x - 5, 162, x + 5, ground_y], fill=(*primary, 210))
        draw.ellipse([x - 24, 122, x + 24, 172], fill=(*primary, 210))

    draw_person(draw, 170, ground_y, primary, accent, scale=1.0)
    draw_bicycle_person(draw, 300, ground_y, primary, accent)
    draw_scooter_person(draw, 690, ground_y, primary, accent)
    img.save(path)
    return path


def draw_person(draw, x: int, ground_y: int, primary: tuple[int, int, int], accent: tuple[int, int, int], scale: float = 1.0) -> None:
    s = scale
    draw.ellipse([x - 8 * s, ground_y - 92 * s, x + 8 * s, ground_y - 76 * s], fill=(*primary, 255))
    draw.line([x, ground_y - 76 * s, x - 12 * s, ground_y - 42 * s], fill=(*accent, 255), width=max(2, int(6 * s)))
    draw.line([x - 7 * s, ground_y - 42 * s, x - 26 * s, ground_y - 4 * s], fill=(*primary, 255), width=max(2, int(5 * s)))
    draw.line([x - 6 * s, ground_y - 42 * s, x + 10 * s, ground_y - 4 * s], fill=(*primary, 255), width=max(2, int(5 * s)))
    draw.line([x - 8 * s, ground_y - 66 * s, x + 20 * s, ground_y - 55 * s], fill=(*primary, 255), width=max(2, int(4 * s)))


def draw_bicycle_person(draw, x: int, ground_y: int, primary: tuple[int, int, int], accent: tuple[int, int, int]) -> None:
    draw.ellipse([x - 70, ground_y - 38, x - 20, ground_y + 12], outline=(*primary, 255), width=5)
    draw.ellipse([x + 25, ground_y - 38, x + 75, ground_y + 12], outline=(*primary, 255), width=5)
    draw.line([x - 45, ground_y - 12, x, ground_y - 50, x + 50, ground_y - 12, x - 45, ground_y - 12], fill=(*primary, 255), width=4)
    draw.line([x, ground_y - 50, x + 20, ground_y - 78], fill=(*primary, 255), width=4)
    draw.ellipse([x + 10, ground_y - 105, x + 28, ground_y - 87], fill=(*primary, 255))
    draw.line([x + 18, ground_y - 88, x - 2, ground_y - 54], fill=(*accent, 255), width=6)
    draw.line([x - 2, ground_y - 54, x + 28, ground_y - 20], fill=(*primary, 255), width=5)


def draw_scooter_person(draw, x: int, ground_y: int, primary: tuple[int, int, int], accent: tuple[int, int, int]) -> None:
    draw.rounded_rectangle([x - 45, ground_y - 42, x + 42, ground_y - 20], radius=10, fill=(*accent, 230))
    draw.ellipse([x - 40, ground_y - 25, x - 18, ground_y - 3], fill=(*primary, 255))
    draw.ellipse([x + 18, ground_y - 25, x + 40, ground_y - 3], fill=(*primary, 255))
    draw.line([x + 35, ground_y - 42, x + 58, ground_y - 86], fill=(*primary, 255), width=4)
    draw.ellipse([x - 8, ground_y - 116, x + 8, ground_y - 100], fill=(*primary, 255))
    draw.line([x, ground_y - 100, x - 2, ground_y - 62], fill=(*accent, 255), width=6)
    draw.line([x - 2, ground_y - 62, x + 32, ground_y - 46], fill=(*primary, 255), width=4)


def rgb_tuple(color: RGBColor) -> tuple[int, int, int]:
    return int(color[0]), int(color[1]), int(color[2])


def highlight_max(chart, points: List[Dict[str, Any]], theme: Dict[str, Any]) -> None:
    try:
        max_idx = max(range(len(points)), key=lambda idx: number(points[idx].get("value")) or 0)
        for idx, point in enumerate(chart.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = theme["accent"] if idx == max_idx else theme["primary"]
    except Exception:
        return


def add_card(slide, left, top, width, height, name: str = "", theme: Dict[str, Any] | None = None, alt: bool = False):
    theme = theme or dashboard_theme({})
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(left), emu(top), emu(width), emu(height))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["card_alt"] if alt else theme["card"]
    shape.line.color.rgb = theme["border"]
    return shape


def add_text(slide, text: str, left, top, width, height, size: int, color=RGBColor(0, 0, 0), bold=False, name: str = "", theme: Dict[str, Any] | None = None):
    theme = theme or dashboard_theme({})
    shape = slide.shapes.add_textbox(emu(left), emu(top), emu(width), emu(height))
    if name:
        shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    run = frame.paragraphs[0].add_run()
    run.text = text
    run.font.name = theme["font"]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def number(value: Any) -> float | None:
    try:
        return float(value)
    except (TypeError, ValueError):
        return None


def emu(value: Any) -> int:
    return int(round(float(value)))


def load_json(path: str | Path) -> Dict[str, Any]:
    return json.loads(Path(path).read_text(encoding="utf-8"))


def run_cli(argv: List[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Create portrait organization dashboard PPTX.")
    parser.add_argument("--package", required=True)
    parser.add_argument("--preflight", required=True)
    parser.add_argument("--template")
    parser.add_argument("--output", required=True)
    args = parser.parse_args(argv)
    output = write_dashboard(args.package, args.preflight, args.output, args.template)
    print(str(output))
    return 0


if __name__ == "__main__":
    raise SystemExit(run_cli())
