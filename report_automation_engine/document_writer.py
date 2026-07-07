"""Create editable PPTX drafts from report_package.json."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Dict, Iterable, List

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ACCENT = RGBColor(0xD9, 0x42, 0x2E)
BLUE = RGBColor(0x2F, 0x66, 0xA7)
GRAY = RGBColor(0x66, 0x66, 0x66)


def write_document(
    package_path: str | Path,
    output_path: str | Path,
    document_type: str,
    preflight_path: str | Path | None = None,
    template_path: str | Path | None = None,
) -> Path:
    package = load_json(package_path)
    preflight = load_json(preflight_path) if preflight_path else {}
    if preflight.get("status") == "blocked":
        raise ValueError("preflight status is blocked; fix errors before creating PPTX output")

    prs = new_presentation(template_path)
    if document_type == "chart_review":
        build_chart_review(prs, package)
    elif document_type == "pptx_report":
        build_report(prs, package)
    else:
        raise ValueError(f"unsupported document type: {document_type}")

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return output


def new_presentation(template_path: str | Path | None) -> Presentation:
    prs = Presentation(str(template_path)) if template_path and Path(template_path).exists() else Presentation()
    # Keep template theme/size, but remove placeholder sample slides from output.
    for slide_id in list(prs.slides._sldIdLst):
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def blank_layout(prs: Presentation):
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]


def build_chart_review(prs: Presentation, package: Dict[str, Any]) -> None:
    add_title_slide(prs, package, "Chart Review Draft")
    charts = [row for row in package.get("charts", []) if row.get("include_chart")]
    if not charts:
        add_text_slide(prs, "검토 필요", ["차트 후보가 없습니다.", "preflight_report.json을 확인하세요."])
        return

    failures = []
    for table_key, rows in group_by(charts, "table_key").items():
        try:
            title = section_title(package, table_key) or table_key or "Chart"
            slide = prs.slides.add_slide(blank_layout(prs))
            add_title(slide, title)
            add_chart(slide, rows, Inches(0.7), Inches(1.35), Inches(7.2), Inches(4.8))
            add_note(slide, chart_note(rows), Inches(8.25), Inches(1.35), Inches(4.3), Inches(4.8))
            add_source(slide, f"source: {table_key}")
        except Exception as exc:
            failures.append(f"{table_key}: {exc}")
    if failures:
        add_text_slide(prs, "검토 필요", failures)


def build_report(prs: Presentation, package: Dict[str, Any]) -> None:
    add_title_slide(prs, package, "Report Draft")
    for section in package.get("sections", []):
        key = str(section.get("table_key", ""))
        slide = prs.slides.add_slide(blank_layout(prs))
        add_title(slide, str(section.get("title") or key or "Untitled"))
        add_textbox(slide, str(section.get("narrative_final", "")), Inches(0.65), Inches(1.05), Inches(5.8), Inches(1.4), 13)
        add_table(slide, table_rows(package, key), Inches(0.65), Inches(2.65), Inches(5.8), Inches(3.7))
        add_chart(slide, chart_rows(package, key), Inches(6.75), Inches(1.35), Inches(5.8), Inches(4.95))
        add_source(slide, f"source: {key}")
    if package.get("qa"):
        add_text_slide(prs, "QA Summary", [f"{q.get('severity', '')}: {q.get('message', '')}" for q in package["qa"][:16]])


def add_title_slide(prs: Presentation, package: Dict[str, Any], fallback_title: str) -> None:
    meta = package.get("meta", {})
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title(slide, str(meta.get("report_title") or meta.get("source_file_name") or fallback_title), top=Inches(2.2), size=30)
    add_textbox(slide, str(meta.get("created_at", "")), Inches(1.0), Inches(3.1), Inches(11.2), Inches(0.5), 14, GRAY)


def add_text_slide(prs: Presentation, title: str, lines: Iterable[str]) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title(slide, title)
    add_textbox(slide, "\n".join(str(line) for line in lines), Inches(0.8), Inches(1.35), Inches(11.6), Inches(5.4), 14)


def add_title(slide, text: str, top=Inches(0.35), size=22) -> None:
    add_textbox(slide, text, Inches(0.65), top, Inches(12.0), Inches(0.6), size, BLUE, bold=True)


def add_textbox(slide, text: str, left, top, width, height, size: int, color=RGBColor(0, 0, 0), bold=False) -> None:
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.word_wrap = True
    frame.clear()
    for idx, line in enumerate(str(text).splitlines() or [""]):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        run = paragraph.add_run()
        run.text = line
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_note(slide, text: str, left, top, width, height) -> None:
    add_textbox(slide, text, left, top, width, height, 12, GRAY)


def add_source(slide, text: str) -> None:
    add_textbox(slide, text, Inches(0.65), Inches(6.85), Inches(12.0), Inches(0.3), 9, GRAY)


def add_table(slide, rows: List[Dict[str, Any]], left, top, width, height) -> None:
    if not rows:
        add_textbox(slide, "삽입표 데이터 없음", left, top, width, height, 11, GRAY)
        return
    visible = rows[:8]
    table_shape = slide.shapes.add_table(len(visible) + 1, 3, left, top, width, height)
    table = table_shape.table
    for idx, heading in enumerate(["항목", "%", "N"]):
        cell = table.cell(0, idx)
        cell.text = heading
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
    for row_idx, row in enumerate(visible, start=1):
        table.cell(row_idx, 0).text = str(row.get("category", ""))
        table.cell(row_idx, 1).text = display_percent(row)
        table.cell(row_idx, 2).text = str(row.get("weighted_n", "") or "")


def add_chart(slide, rows: List[Dict[str, Any]], left, top, width, height) -> None:
    rows = [row for row in rows if number(row.get("value")) is not None]
    if not rows:
        add_textbox(slide, "차트 데이터 없음", left, top, width, height, 12, GRAY)
        return
    rows = sorted(rows, key=lambda row: number(row.get("sort_order")) or 9999)[:10]
    chart_data = ChartData()
    chart_data.categories = [str(row.get("category", "")) for row in rows]
    chart_data.add_series("값", [number(row.get("value")) or 0 for row in rows])
    chart_type = XL_CHART_TYPE.PIE if choose_chart_type(rows) == "pie" else XL_CHART_TYPE.COLUMN_CLUSTERED
    chart = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data).chart
    chart.has_legend = chart_type == XL_CHART_TYPE.PIE
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
    highlight_max(chart, rows)


def highlight_max(chart, rows: List[Dict[str, Any]]) -> None:
    max_idx = max(range(len(rows)), key=lambda idx: number(rows[idx].get("value")) or 0)
    points = chart.series[0].points
    for idx, point in enumerate(points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = ACCENT if idx == max_idx else BLUE


def chart_note(rows: List[Dict[str, Any]]) -> str:
    chart_type = choose_chart_type(rows)
    highlight = max(rows, key=lambda row: number(row.get("value")) or 0)
    lines = [
        f"차트 유형: {chart_type}",
        f"강조 항목: {highlight.get('category', '')} ({highlight.get('display_value') or highlight.get('value', '')})",
        "",
        "데이터",
    ]
    lines.extend(f"- {row.get('category', '')}: {row.get('display_value') or row.get('value', '')}" for row in rows[:10])
    return "\n".join(lines)


def table_rows(package: Dict[str, Any], table_key: str) -> List[Dict[str, Any]]:
    for table in package.get("tables", []):
        if table.get("table_key") == table_key:
            return list(table.get("rows", []))
    return []


def chart_rows(package: Dict[str, Any], table_key: str) -> List[Dict[str, Any]]:
    return [row for row in package.get("charts", []) if row.get("table_key") == table_key and row.get("include_chart")]


def section_title(package: Dict[str, Any], table_key: str) -> str:
    for section in package.get("sections", []):
        if section.get("table_key") == table_key:
            return str(section.get("title") or "")
    return ""


def choose_chart_type(rows: List[Dict[str, Any]]) -> str:
    categories = {str(row.get("category", "")) for row in rows if row.get("category")}
    return "pie" if len(categories) == 2 else "column"


def display_percent(row: Dict[str, Any]) -> str:
    value = row.get("percent")
    unit = row.get("unit") or "%"
    return "" if value in (None, "") else f"{value}{unit}"


def group_by(rows: List[Dict[str, Any]], key: str) -> Dict[str, List[Dict[str, Any]]]:
    grouped: Dict[str, List[Dict[str, Any]]] = {}
    for row in rows:
        grouped.setdefault(str(row.get(key, "")), []).append(row)
    return grouped


def number(value: Any) -> float | None:
    try:
        return float(value)
    except (TypeError, ValueError):
        return None


def load_json(path: str | Path | None) -> Dict[str, Any]:
    if not path:
        return {}
    return json.loads(Path(path).read_text(encoding="utf-8"))


def run_cli(argv: List[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Create editable PPTX drafts from report_package.json.")
    parser.add_argument("--package", required=True)
    parser.add_argument("--preflight")
    parser.add_argument("--template")
    parser.add_argument("--type", required=True, choices=["pptx_report", "chart_review"])
    parser.add_argument("--output", required=True)
    args = parser.parse_args(argv)
    output = write_document(args.package, args.output, args.type, args.preflight, args.template)
    print(str(output))
    return 0


if __name__ == "__main__":
    raise SystemExit(run_cli())
